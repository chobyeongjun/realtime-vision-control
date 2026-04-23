"""
H-Walker 2026-05-04 Meeting PPT Builder
Takes VRE_Template.pptx, populates it with real experiment data, reorders, and saves.
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from copy import deepcopy
from pathlib import Path

BASE = Path("/sessions/beautiful-dazzling-tesla/mnt/realtime-vision-control")
TEMPLATE = BASE / ".templates" / "template.pptx"
FIG_DIR  = BASE / "docs" / "meetings" / "2026-05-04" / "figures"
OUT_PPTX = BASE / "docs" / "meetings" / "2026-05-04" / "presentation.pptx"


# ============================================================
# Helpers
# ============================================================
NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'


def set_text(shape, new_text):
    """Replace text while *exactly* preserving first run's XML (font, Asian typeface, color).

    Korean characters rely on <a:rPr><a:ea typeface="..."/></a:rPr>. Copying only the
    python-pptx `font.name` loses `ea` and `cs` typefaces → tofu/garbage glyphs. We
    deep-copy the whole <a:rPr>/<a:pPr> XML and replace only the text node.
    """
    if not shape.has_text_frame:
        return False
    from lxml import etree
    body = shape.text_frame._txBody

    ps = body.findall(f'{NS}p')
    if not ps:
        return False
    first_p = ps[0]

    # Grab template <a:pPr> and <a:rPr> (if present) from first paragraph / first run
    template_pPr = first_p.find(f'{NS}pPr')
    template_r   = first_p.find(f'{NS}r')
    template_rPr = template_r.find(f'{NS}rPr') if template_r is not None else None

    # Remove all existing <a:p>
    for p in ps:
        body.remove(p)

    lines = new_text.split("\n")
    for line in lines:
        p_el = etree.SubElement(body, f'{NS}p')
        if template_pPr is not None:
            p_el.append(deepcopy(template_pPr))
        r_el = etree.SubElement(p_el, f'{NS}r')
        rPr_for_line = deepcopy(template_rPr) if template_rPr is not None else etree.SubElement(
            etree.Element(f'{NS}tmp'), f'{NS}rPr'
        )
        if template_rPr is None:
            # Built a fresh rPr; attach manually
            r_el.append(rPr_for_line)
        else:
            r_el.append(rPr_for_line)
        # Safety net: if line contains Korean, guarantee East-Asian typeface present
        has_korean = any(0xAC00 <= ord(c) <= 0xD7A3 for c in line)
        if has_korean and rPr_for_line.find(f'{NS}ea') is None:
            ea_el = etree.SubElement(rPr_for_line, f'{NS}ea')
            ea_el.set('typeface', 'Pretendard')
        t_el = etree.SubElement(r_el, f'{NS}t')
        t_el.text = line
    return True


def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_slide_text(slide, name, text):
    sh = find_shape(slide, name)
    if sh is None:
        print(f"  [WARN] shape '{name}' not found")
        return False
    return set_text(sh, text)


def add_picture_over_shape(slide, placeholder_name, image_path):
    """Find a placeholder shape, get its geometry, remove it, add picture in same position."""
    sh = find_shape(slide, placeholder_name)
    if sh is None:
        # try to find FIGURE/IMAGE text placeholder
        return None
    left, top, width, height = sh.left, sh.top, sh.width, sh.height
    # Insert picture (keep placeholder below for border)
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    return pic


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def reorder_and_keep_slides(prs, keep_indices_in_order):
    """
    Rearrange presentation slides to the given order and delete the rest.
    keep_indices_in_order: list of 0-based indices from the original template.
    """
    xml_slides = prs.slides._sldIdLst  # CT_SlideIdList
    slides_list = list(xml_slides)
    # Build new order
    new_list = [slides_list[i] for i in keep_indices_in_order]
    # Remove all
    for s in slides_list:
        xml_slides.remove(s)
    # Add in new order
    for s in new_list:
        xml_slides.append(s)


# ============================================================
# Build
# ============================================================
def main():
    prs = Presentation(str(TEMPLATE))

    # =====================================================
    # Slide 1 (template index 0) — Cover
    # =====================================================
    s = prs.slides[0]
    set_slide_text(s, "MainTitle", "H-Walker\nReal-Time Vision Control")
    set_slide_text(s, "Subtitle", "케이블 드리븐 보행 보조기를 위한 실시간 하체 포즈 추정 시스템")
    set_slide_text(s, "AuthorInfo", "조병준  |  중앙대학교 기계공학과  |  AR Lab")
    set_slide_text(s, "Tag", "FIRST MEETING")

    # =====================================================
    # Slide 2 (index 1) — AGENDA
    # =====================================================
    s = prs.slides[1]
    set_slide_text(s, "Text 12", "Introduction")
    set_slide_text(s, "Text 13", "실시간 비전 제어의 동기와 목표")
    set_slide_text(s, "Text 16", "Experiments")
    set_slide_text(s, "Text 17", "4단계 Evolution · 44.4 → 13.7 ms")
    set_slide_text(s, "Text 20", "Discussion")
    set_slide_text(s, "Text 21", "왜 빨라졌는가 · 실패 교훈 · 방향성 질문")
    set_slide_text(s, "Text 24", "Next Steps")
    set_slide_text(s, "Text 25", "다음 2주 Task")

    # =====================================================
    # Slide 3 (index 2) — Motivation: 왜 '실시간 비전 제어'인가 (3 bullets)
    # =====================================================
    s = prs.slides[2]
    set_slide_text(s, "Text 5", "연구 동기 — 왜 실시간 비전 제어인가")
    set_slide_text(s, "Text 6", "보행 보조 제어의 핵심은 사용자 의도 측정. 무엇으로, 얼마나 빠르게, 얼마나 안전하게 받아오느냐가 결정한다.")
    set_slide_text(s, "Text 9", "사용자 의도의 직접 관측")
    set_slide_text(s, "Text 10", "힘센서·버튼은 결과만 감지. 관절 포즈는 보행 위상·의도 그 자체.")
    set_slide_text(s, "Text 13", "임피던스 제어의 50ms 벽")
    set_slide_text(s, "Text 14", "Human reaction ≒ 100ms. 센싱이 그 절반 안에 들어와야 능동 보조가 안정화.")
    set_slide_text(s, "Text 17", "마커리스·비접촉 요건")
    set_slide_text(s, "Text 18", "임상 적용 시 마커 부착·Wearable 장착 불가. 단일 카메라로 해결해야 함.")
    set_slide_text(s, "Text 22", "Motivation")
    set_slide_text(s, "Text 23", "Real-Time Vision Control")
    set_slide_text(s, "Text 24", "Fig. 1 — 사용자 의도를 관절 포즈로 추정하는 비전 제어 루프")

    # =====================================================
    # Slide 4 (index 3) — Two-column: 포즈 추정 접근법 비교
    # =====================================================
    s = prs.slides[3]
    set_slide_text(s, "Text 5", "포즈 추정 접근법 — 왜 Vision을 택했는가")
    set_slide_text(s, "Text 7", "기존 접근의 한계")
    set_slide_text(s, "Text 10", "Marker-MoCap")
    set_slide_text(s, "Text 11", "정확도 최고, 그러나 마커 부착·다카메라 고정 공간 필요. 임상 불가.")
    set_slide_text(s, "Text 13", "Wearable IMU")
    set_slide_text(s, "Text 14", "Drift 누적·장착 재현성 문제. 환자가 매번 착용해야 함.")
    set_slide_text(s, "Text 16", "Contact Force/EMG")
    set_slide_text(s, "Text 17", "결과만 측정, 의도·위상 직접 관측 불가.")
    set_slide_text(s, "Text 19", "제안 — 단일 카메라 Vision")
    set_slide_text(s, "Text 22", "비접촉·마커리스")
    set_slide_text(s, "Text 23", "ZED X Mini 1대. 환자 장착 0, 공간 제약 최소.")
    set_slide_text(s, "Text 25", "관절 포즈 직접 추정")
    set_slide_text(s, "Text 26", "하체 6 kpt 3D. 의도·위상·자세 한 번에.")
    set_slide_text(s, "Text 28", "실시간·안전 보장")
    set_slide_text(s, "Text 29", "E2E 13.7 ms, 20 ms HARD LIMIT, violation 0.")

    # =====================================================
    # Slide 5 (index 4) — 연구 목표 (간결)
    # =====================================================
    s = prs.slides[4]
    set_slide_text(s, "Text 5", "연구 목표 — 실시간 비전 기반 제어 루프")
    set_slide_text(s, "Text 7", "OBJECTIVE")
    set_slide_text(s, "Text 8", "하체 6 관절 3D 포즈를 E2E < 50 ms 로 공급하여 임피던스 제어의 관측값으로 사용.")
    set_slide_text(s, "Text 9", "CORE STACK")
    set_slide_text(s, "Text 10", "ZED X Mini · YOLO26s-lower6 · TRT FP16 · 20 ms HARD LIMIT")
    set_slide_text(s, "Text 11", "OUTCOME")
    set_slide_text(s, "Text 12", "E2E 13.7 ms · p99 19.8 ms · motor violation 0건")
    set_slide_text(s, "Text 22", "Pipeline")
    set_slide_text(s, "Text 23", "Camera → Pose → 3D → SHM → Control")
    set_slide_text(s, "Text 24", "Fig. 2 — 전체 파이프라인 개요")

    # =====================================================
    # Slide 7 (index 6) — DELETE (process flow T6 — not needed here)
    # (keep slide index 5 = T6) — we'll DELETE this one
    # =====================================================

    # =====================================================
    # Slide 7 (index 6) — Comparison table T7: 실험 1 모델 벤치마크
    # =====================================================
    s = prs.slides[6]
    set_slide_text(s, "Title", "실험 1 — 12개 Pose 모델 벤치마크")
    set_slide_text(s, "Header0", "지표")
    set_slide_text(s, "Header1", "YOLOv8n-Pose  ★")
    set_slide_text(s, "Header2", "YOLO26s-Pose")
    set_slide_text(s, "Header3", "MediaPipe / RTMPose")
    set_slide_text(s, "Cell0_0", "E2E mean [ms]")
    set_slide_text(s, "Cell0_1", "35.7  (가장 빠름)")
    set_slide_text(s, "Cell0_2", "44.4  (<50ms 아슬아슬)")
    set_slide_text(s, "Cell0_3", "63.6  /  150.5+  (실패)")
    set_slide_text(s, "Cell1_0", "FPS")
    set_slide_text(s, "Cell1_1", "27.7")
    set_slide_text(s, "Cell1_2", "22.3")
    set_slide_text(s, "Cell1_3", "15.5 / 6.6")
    set_slide_text(s, "Cell2_0", "E2E <50ms 비율")
    set_slide_text(s, "Cell2_1", "100 %")
    set_slide_text(s, "Cell2_2", "99.4 %")
    set_slide_text(s, "Cell2_3", "0 %")
    set_slide_text(s, "Cell3_0", "인식률 / Confidence")
    set_slide_text(s, "Cell3_1", "100 % / 0.97")
    set_slide_text(s, "Cell3_2", "100 % / 0.99")
    set_slide_text(s, "Cell3_3", "94.4 % / 0.77")
    set_slide_text(s, "Footnote", "* Jetson Orin NX 16GB · ZED X Mini SVGA@120 · TRT FP16 · 15s 측정")

    # =====================================================
    # Slide 8 (index 7) — T8 3 KPIs: 실험 4 Hard Limit
    # =====================================================
    s = prs.slides[7]
    set_slide_text(s, "Text 5", "실험 4 — 20 ms HARD LIMIT 안전 구조")
    set_slide_text(s, "Text 6", "300초 연속 운전 · frame-skip 패턴 기반 정량적 안전 보장")
    set_slide_text(s, "Text 9", "19.8 ms")
    set_slide_text(s, "Text 10", "p99 E2E Latency")
    set_slide_text(s, "Text 11", "HARD LIMIT 20 ms 기준")
    set_slide_text(s, "Text 14", "0.031 %")
    set_slide_text(s, "Text 15", "HARD LIMIT 초과율")
    set_slide_text(s, "Text 16", "93 / 300,000 frame")
    set_slide_text(s, "Text 19", "0 건")
    set_slide_text(s, "Text 20", "모터 도달 violation")
    set_slide_text(s, "Text 21", "valid=False 로 제어 차단")
    set_slide_text(s, "Text 22", "* 5-layer defense: GC off · SCHED_FIFO 90 · taskset · watchdog · 70 N clamp")

    # =====================================================
    # Slide 9 (index 8) — T9 4 KPIs: 종합 성과
    # =====================================================
    s = prs.slides[8]
    set_slide_text(s, "직사각형 1", "종합 성과 — 3.24× 가속 + 안전 보장")
    set_slide_text(s, "TextBox 2", "Baseline 44.4 ms · 22 FPS → 최종 13.7 ms · 73 FPS · violation 0")
    set_slide_text(s, "TextBox 6", "3.24×")
    set_slide_text(s, "TextBox 7", "E2E Speedup")
    set_slide_text(s, "TextBox 8", "44.4 → 13.7 ms")
    set_slide_text(s, "TextBox 11", "13.7 ms")
    set_slide_text(s, "TextBox 12", "Final E2E mean")
    set_slide_text(s, "TextBox 13", "Jetson 실측")
    set_slide_text(s, "TextBox 16", "73 FPS")
    set_slide_text(s, "TextBox 17", "Pose 처리율")
    set_slide_text(s, "TextBox 18", "Python 단독")
    set_slide_text(s, "TextBox 21", "0 건")
    set_slide_text(s, "TextBox 22", "Motor violation")
    set_slide_text(s, "TextBox 23", "20 ms HARD LIMIT")

    # =====================================================
    # Slide 11 (index 10) — T11 big image: 시스템 아키텍처
    # =====================================================
    s = prs.slides[10]
    set_slide_text(s, "직사각형 1", "시스템 아키텍처")
    set_slide_text(s, "TextBox 2", "Perception → IPC · Safety → C++ Impedance → Motor")
    set_slide_text(s, "TextBox 9", "Fig. 3 — 전체 파이프라인 (E2E 13.7 ms, p99 19.8 ms)")
    # Insert image
    ph_tl = find_shape(s, "직사각형 4")  # image background
    if ph_tl:
        left, top, width, height = ph_tl.left, ph_tl.top, ph_tl.width, ph_tl.height
        # Remove placeholder texts
        for n in ["TextBox 5", "TextBox 6"]:
            p = find_shape(s, n)
            if p:
                remove_shape(p)
        s.shapes.add_picture(str(FIG_DIR / "fig_architecture.png"),
                             left, top, width=width, height=height)

    # =====================================================
    # Slide (index 9) — T10 quote: ★ 왜 빨라졌는가 (6 카테고리 철학)
    # Shapes: 직사각형 1 (title) / TextBox 2 (subtitle) /
    #         TextBox 6 (opening quote mark) / TextBox 7 (quote body) /
    #         TextBox 8 (author) / TextBox 9 (footnote)
    # =====================================================
    s = prs.slides[9]
    set_slide_text(s, "직사각형 1", "왜 빨라졌는가 — 무엇을 정상화했는가")
    set_slide_text(s, "TextBox 2", "44.4 → 13.7 ms 가속의 원인을 6개 카테고리로 분해")
    set_slide_text(s, "TextBox 7",
        "1. 연산량 감소 (-26.4 ms) : Head 17→6 kpt\n"
        "2. 파이프라인 오버랩 (-9.4 ms) : ready_rgb / depth split\n"
        "3. 프레임워크 우회 (-9 ms) : DirectTRT + C++ post\n"
        "4. GPU·OS 최대화 : jetson_clocks + MAXN + gc.off\n"
        "5. 경합 제거 : CPU isolation + SCHED_FIFO\n"
        "6. 실시간성 보장 : 20 ms HARD LIMIT + SHM seqlock"
    )
    set_slide_text(s, "TextBox 8", "— 핵심 설계 원칙")
    set_slide_text(s, "TextBox 9",
        "\"속도는 자원 분배의 문제\" — I/O는 숨기고, 프레임워크는 벗기고, 하드웨어는 고정."
    )

    # =====================================================
    # Slide (index 11) — T12 3-step: 실패 교훈 (무엇을 하지 말아야 하는가)
    # Shapes (title / body pairs):
    #   Step 1 = TextBox 7  / TextBox 9
    #   Step 2 = TextBox 13 / TextBox 15
    #   Step 3 = TextBox 19 / TextBox 21
    # =====================================================
    s = prs.slides[11]
    set_slide_text(s, "직사각형 1", "실패 교훈 — 역효과였던 접근 3가지")
    set_slide_text(s, "TextBox 2", "논문 Methods 섹션의 결정적 근거로 남음")

    mapping_12 = {
        "TextBox 7":  "2D 필터링 금지",
        "TextBox 9":
            "One Euro / SegLen 을 2D에 적용 시\n"
            "depth NaN → 3D 실패.\n"
            "규칙: 필터는 3D 단계에서만.",
        "TextBox 13": "NEURAL Depth 기각",
        "TextBox 15":
            "2 cm 정확도 vs GPU SM 경합.\n"
            "predict 2.4× 급락, 29 FPS.\n"
            "규칙: PERFORMANCE 유지.",
        "TextBox 19": "Zero-Copy 포기",
        "TextBox 21":
            "copy=False → capture overwrite race.\n"
            "calib 100 → 0 %.\n"
            "규칙: 공유 버퍼는 복사 필수.",
    }
    for name, val in mapping_12.items():
        set_slide_text(s, name, val)

    # =====================================================
    # Slide 13 (index 12) — T13 4-feature: ★ 방향성 질문
    # =====================================================
    s = prs.slides[12]
    set_slide_text(s, "직사각형 1", "교수님 가이드 요청 — 방향성 열린 질문")
    set_slide_text(s, "TextBox 2", "앞으로 2~3개월 방향성 설정에 도움을 구하는 4가지")
    set_slide_text(s, "TextBox 5", "논문 타깃")
    set_slide_text(s, "TextBox 7", "RA-L / IROS·ICRA /\n임상 저널 중\n어디를 우선?")
    set_slide_text(s, "TextBox 9", "비전-IMU 융합 깊이")
    set_slide_text(s, "TextBox 11", "Static R 유지?\n완전 VIO 까지?")
    set_slide_text(s, "TextBox 13", "Sim-to-Real RL 시점")
    set_slide_text(s, "TextBox 15", "지금 착수 vs\n제어 검증 후 착수?")
    set_slide_text(s, "TextBox 17", "Clinical · HW 방향")
    set_slide_text(s, "TextBox 19", "IRB 진입 시점?\n케이블 vs 직구동?")

    # =====================================================
    # Slide 14 (index 13) — T14 1 photo: Latency bar chart
    # =====================================================
    s = prs.slides[13]
    set_slide_text(s, "Title", "실험 1 — E2E Latency 비교")
    set_slide_text(s, "Subtitle", "YOLO 계열 6종만 50 ms 벽 통과. MediaPipe · RTMPose 실패.")
    set_slide_text(s, "Caption", "Fig. 4 — 12개 Pose 모델 E2E Latency")
    # Insert image
    imgbg = find_shape(s, "ImageBg")
    if imgbg:
        left, top, width, height = imgbg.left, imgbg.top, imgbg.width, imgbg.height
        # Remove ImgLabel placeholder
        p = find_shape(s, "ImgLabel")
        if p: remove_shape(p)
        s.shapes.add_picture(str(FIG_DIR / "fig_latency_bar.png"),
                             left, top, width=width, height=height)

    # =====================================================
    # Slide 15 (index 14) — T15 2 photos: Fine-Tuning + Evolution
    # =====================================================
    s = prs.slides[14]
    set_slide_text(s, "Title", "실험 2 · 3 — Fine-Tuning + Pipeline 최적화")
    set_slide_text(s, "Subtitle", "6 kpt head 전환 + DirectTRT + C++ post → 44 ms에서 13.7 ms까지.")
    set_slide_text(s, "Caption1", "Fig. 5 — Fine-Tuning 4대 지표 (17 → 6 kpt)")
    set_slide_text(s, "Caption2", "Fig. 6 — 4단계 Evolution & 기여도")
    # Insert two images
    imgbg1 = find_shape(s, "ImgBg210")
    if imgbg1:
        left, top, width, height = imgbg1.left, imgbg1.top, imgbg1.width, imgbg1.height
        p = find_shape(s, "ImgLabel_210")
        if p: remove_shape(p)
        s.shapes.add_picture(str(FIG_DIR / "fig_ft_compare.png"),
                             left, top, width=width, height=height)
    imgbg2 = find_shape(s, "ImgBg250")
    if imgbg2:
        left, top, width, height = imgbg2.left, imgbg2.top, imgbg2.width, imgbg2.height
        p = find_shape(s, "ImgLabel_250")
        if p: remove_shape(p)
        s.shapes.add_picture(str(FIG_DIR / "fig_evolution.png"),
                             left, top, width=width, height=height)

    # =====================================================
    # Slide 18 (index 17) — T18 timeline: 실험 로드맵 (with "why" per phase)
    # =====================================================
    s = prs.slides[17]
    set_slide_text(s, "직사각형 1", "실험 로드맵 — 4단계 Evolution")
    set_slide_text(s, "TextBox 2", "44.4 ms → 13.7 ms (3.24×), motor violation 0")
    set_slide_text(s, "TextBox 6", "Phase 1")
    set_slide_text(s, "TextBox 7", "모델 선정")
    set_slide_text(s, "TextBox 8",
        "12개 모델 벤치마크.\nYOLO만 50 ms 통과.")
    set_slide_text(s, "TextBox 10", "Phase 2")
    set_slide_text(s, "TextBox 11", "Fine-Tuning")
    set_slide_text(s, "TextBox 12",
        "6 kpt head.\n44 → 18 ms (2.4×).\nCat 1: 연산량.")
    set_slide_text(s, "TextBox 14", "Phase 3")
    set_slide_text(s, "TextBox 15", "Pipeline")
    set_slide_text(s, "TextBox 16",
        "DirectTRT + C++ post.\n18 → 13.7 ms / 73 FPS.\nCat 2 · 3.")
    set_slide_text(s, "TextBox 18", "Phase 4")
    set_slide_text(s, "TextBox 19", "Safety")
    set_slide_text(s, "TextBox 20",
        "p99 19.8 ms.\nviolation 0.\nCat 4 · 5 · 6.")

    # =====================================================
    # Slide 6 (index 5) — T6 process flow: 다음 2주 Task
    # (keeping this one — but will repurpose for Next 2 weeks plan)
    # =====================================================
    s = prs.slides[5]
    set_slide_text(s, "직사각형 1", "다음 2주 Task")
    set_slide_text(s, "TextBox 2", "Paper-quality 수치 확보 · 실모터 구동 dry-run")
    # T6 flow: (step body = TextBox 7/12/17/22, step label = TextBox 23/24/25/26)
    mapping_6 = {
        "TextBox 7":  "Baseline 재측정\np95 / std 기록",
        "TextBox 12": "좌우 비대칭\n원인 규명",
        "TextBox 17": "Teensy + AK60\n70N 구동",
        "TextBox 22": "미팅 질문 정리\n+ 발표 리허설",
        "TextBox 23": "재측정",
        "TextBox 24": "캘리브",
        "TextBox 25": "실모터",
        "TextBox 26": "리허설",
    }
    for n, v in mapping_6.items():
        set_slide_text(s, n, v)

    # =====================================================
    # Slide 19 (index 18) — T19 Thanks (keep as-is)
    # =====================================================
    s = prs.slides[18]
    set_slide_text(s, "Text 5", "AR Lab  ·  Assistive & Rehabilitation Robotics Lab  ·  Chung-Ang University  ·  조병준")

    # =====================================================
    # Reorder & Delete unused
    # =====================================================
    # Desired output order (0-based template indices):
    #  1 Cover (0)
    #  2 AGENDA (1)
    #  3 Motivation (2)
    #  4 Passive vs H-Walker (3)
    #  5 Goal & System (4)
    #  6 Architecture image (10)
    #  7 Roadmap timeline (17)        ← Category별 기여 표시
    #  8 Benchmark table (6)          — 실험 1
    #  9 Latency bar chart (13)       — 실험 1 시각화
    # 10 Fine-Tuning + Evolution (14) — 실험 2+3
    # 11 ★ Why-it-got-faster quote (9) — 6 카테고리 철학       [NEW]
    # 12 실패 교훈 (11)              — 3 failed approaches    [NEW]
    # 13 Hard Limit KPI (7)          — 실험 4
    # 14 Final KPI (8)               — 종합
    # 15 Direction questions (12)    — ★ 열린 질문
    # 16 Next 2 weeks (5)
    # 17 Thanks (18)
    # DELETE: T16 3-photo (15), T17 4-photo (16)

    desired = [0, 1, 2, 3, 4, 10, 17, 6, 13, 14, 9, 11, 7, 8, 12, 5, 18]
    reorder_and_keep_slides(prs, desired)

    # Save
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"Saved: {OUT_PPTX}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
